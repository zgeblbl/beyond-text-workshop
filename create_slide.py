from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

# --- AYARLAR VE RENKLER ---
PRS_WIDTH = Inches(13.333) # Geniş ekran (16:9)
PRS_HEIGHT = Inches(7.5)
BLUE_CV = RGBColor(41, 128, 185)
GREEN_AI = RGBColor(39, 174, 96)
ORANGE_UI = RGBColor(230, 126, 34)
DARK_GRAY = RGBColor(50, 50, 50)

prs = Presentation()
prs.slide_width = PRS_WIDTH
prs.slide_height = PRS_HEIGHT

# --- YARDIMCI FONKSİYONLAR ---
def add_title_slide(prs, title_text, subtitle_text):
    slide_layout = prs.slide_layouts[0] # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = title_text
    title.text_frame.paragraphs[0].font.color.rgb = DARK_GRAY
    title.text_frame.paragraphs[0].font.bold = True
    
    subtitle.text = subtitle_text
    subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)

def add_content_slide(prs, title_text):
    slide_layout = prs.slide_layouts[5] # Title Only
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title_text
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = DARK_GRAY
    return slide

def create_styled_box(slide, left, top, width, height, text, color, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = RGBColor(255, 255, 255)
    
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shape

def connect_shapes(slide, shape1, shape2):
    connector = slide.shapes.add_connector(MSO_CONNECTOR.ELBOW, 0, 0, 0, 0)
    connector.begin_connect(shape1, 3) # 3: Sağ kenar
    connector.end_connect(shape2, 1)   # 1: Sol kenar
    connector.line.width = Pt(3)
    connector.line.color.rgb = DARK_GRAY

# ==============================================================================
# SLAYT 1: BAŞLIK
# ==============================================================================
add_title_slide(prs, "BEYOND TEXT", "Görüntü İşleme ve Üretken Yapay Zeka ile\nMultimodal Etkileşim Workshop'u")

# ==============================================================================
# SLAYT 2: BÜYÜK RESİM (AKIŞ DİYAGRAMI)
# ==============================================================================
slide2 = add_content_slide(prs, "Sistemin Akışı: Veri Nasıl Dolaşıyor?")

# Kutular
box_cam = create_styled_box(slide2, Inches(0.5), Inches(3), Inches(2), Inches(1.5), "📷\nKamera\n(Giriş)", DARK_GRAY)
box_cv = create_styled_box(slide2, Inches(3.5), Inches(3), Inches(2.5), Inches(1.5), "👁️ Gözler\n(MediaPipe)\n[Koordinatlar]", BLUE_CV)
box_ai = create_styled_box(slide2, Inches(7), Inches(3), Inches(2.5), Inches(1.5), "🧠 Beyin\n(LLM/Gemini)\n[Yorumlama]", GREEN_AI)
box_ui = create_styled_box(slide2, Inches(10.5), Inches(3), Inches(2), Inches(1.5), "💻\nGradio\n(Çıkış)", ORANGE_UI)

# Oklar
connect_shapes(slide2, box_cam, box_cv)
connect_shapes(slide2, box_cv, box_ai)
connect_shapes(slide2, box_ai, box_ui)

# ==============================================================================
# SLAYT 3: DERİN DALIŞ - MEDIAPIPE (GÖRSELLİ!)
# ==============================================================================
slide3 = add_content_slide(prs, "Gözler: Google MediaPipe Nasıl Görüyor?")

# Sol Taraf: Açıklama Metni
left_textbox = slide3.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5), Inches(5))
tf = left_textbox.text_frame
tf.word_wrap = True
p = tf.add_paragraph()
p.text = "Bilgisayar 'El' Görmez, 'Nokta' Görür."
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = BLUE_CV

bullets = [
    "MediaPipe, el üzerinde 21 adet kritik nokta (Landmark) tespit eder.",
    "Her noktanın bir ID'si ve 3D koordinatı (x, y, z) vardır.",
    "Örnek: ID 8 = İşaret Parmağı Ucu",
    "Biz bu koordinatları kıyaslayarak (if/else) hareketi anlarız."
]
for bullet in bullets:
    p = tf.add_paragraph()
    p.text = bullet
    p.font.size = Pt(18)
    p.level = 1

# Sağ Taraf: GÖRSEL EKLEME (En önemli kısım)
img_path = 'wp.jpg' # <-- BU DOSYA KODLA AYNI YERDE OLMALI
try:
    # Resmi ekle ve sağa hizala
    slide3.shapes.add_picture(img_path, Inches(6), Inches(1.5), width=Inches(6.5))
except FileNotFoundError:
    # Eğer resim yoksa uyarı kutusu koy
    err_box = create_styled_box(slide3, Inches(6), Inches(2), Inches(6), Inches(3), 
                                "⚠️ GÖRSEL EKSİK!\n'hand_landmarks.png' dosyasını\nscriptin yanına koymadınız.", RGBColor(231, 76, 60))

# ==============================================================================
# SLAYT 4: DERİN DALIŞ - LLM (YAPAY ZEKA)
# ==============================================================================
slide4 = add_content_slide(prs, "Beyin: LLM'e Rol Biçmek (Prompt Engineering)")

# Sol Kutu: Standart Kod
box_std = create_styled_box(slide4, Inches(1), Inches(2), Inches(5), Inches(4), "", DARK_GRAY, MSO_SHAPE.RECTANGLE)
tf = box_std.text_frame
p = tf.paragraphs[0]
p.text = "Eski Usul Programlama"
p.font.size = Pt(24)
p = tf.add_paragraph()
p.text = "\nif gesture == 'Zafer':\n    print('Barış işareti yapıldı.')"
p.font.name = 'Courier New'
p.font.size = Pt(16)
p.alignment = PP_ALIGN.LEFT

# Sağ Kutu: AI Yaklaşımı
box_ai_deep = create_styled_box(slide4, Inches(7), Inches(2), Inches(5), Inches(4), "", GREEN_AI, MSO_SHAPE.RECTANGLE)
tf = box_ai_deep.text_frame
p = tf.paragraphs[0]
p.text = "Yapay Zeka Yaklaşımı"
p.font.size = Pt(24)
p = tf.add_paragraph()
p.text = "\nPROMPT (Kimlik):\n'Sen huysuz bir korsan robotsun.'\n\nUSER GİRDİSİ:\n'Kullanıcı Zafer işareti yaptı.'\n\nAI ÇIKTISI:\n'Arrgh! Barış mı? Denizlerde barış olmaz evlat!'"
p.font.name = 'Calibri'
p.font.size = Pt(18)
p.alignment = PP_ALIGN.LEFT

# ==============================================================================
# KAYDET
# ==============================================================================
output_filename = 'Workshop_Sunumu_Pro.pptx'
prs.save(output_filename)
print(f"✅ Profesyonel sunum hazırlandı: {output_filename}")
print("👉 Google Slides'a yükleyip kullanabilirsin!")